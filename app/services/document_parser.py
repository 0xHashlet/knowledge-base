from dataclasses import dataclass


@dataclass(frozen=True)
class ParsedDocument:
    text: str
    parser_name: str
    parser_version: str


_MIME_PARSER_MAP = {
    "text/plain": ("plain_text", "_parse_text"),
    "text/markdown": ("markdown", "_parse_text"),
    "application/pdf": ("pymupdf", "_parse_pdf"),
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": (
        "python-docx",
        "_parse_docx",
    ),
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": (
        "openpyxl",
        "_parse_xlsx",
    ),
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": (
        "python-pptx",
        "_parse_pptx",
    ),
}


class DocumentParser:
    parser_version = "2"

    def parse(self, *, file_type: str, content: bytes) -> ParsedDocument:
        if file_type not in _MIME_PARSER_MAP:
            raise ValueError(f"Unsupported document file type: {file_type}")
        parser_name, method_name = _MIME_PARSER_MAP[file_type]
        parsed_text = getattr(self, method_name)(content)
        return ParsedDocument(text=parsed_text, parser_name=parser_name, parser_version=self.parser_version)

    @staticmethod
    def _parse_text(content: bytes) -> str:
        return content.decode("utf-8")

    @staticmethod
    def _parse_pdf(content: bytes) -> str:
        import fitz

        with fitz.open(stream=content, filetype="pdf") as pdf:
            return "\n".join(page.get_text() for page in pdf)

    @staticmethod
    def _parse_docx(content: bytes) -> str:
        from docx import Document
        from io import BytesIO

        doc = Document(BytesIO(content))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        if not paragraphs:
            paragraphs = [p.text for p in doc.paragraphs]
        return "\n".join(paragraphs)

    @staticmethod
    def _parse_xlsx(content: bytes) -> str:
        from io import BytesIO
        import openpyxl

        wb = openpyxl.load_workbook(BytesIO(content), read_only=True, data_only=True)
        parts = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            parts.append(f"[工作表: {sheet_name}]")
            rows = list(ws.iter_rows(values_only=True))
            if not rows:
                continue
            headers = [str(h) if h is not None else "" for h in rows[0]]
            for row in rows[1:]:
                values = [str(v) if v is not None else "" for v in row]
                pairs = [f"{h}: {v}" for h, v in zip(headers, values) if h or v]
                if pairs:
                    parts.append("; ".join(pairs))
        wb.close()
        return "\n".join(parts)

    @staticmethod
    def _parse_pptx(content: bytes) -> str:
        from io import BytesIO
        from pptx import Presentation

        prs = Presentation(BytesIO(content))
        parts = []
        for i, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            texts.append(text)
            if texts:
                parts.append(f"[幻灯片 {i}]\n" + "\n".join(texts))
        return "\n\n".join(parts)
